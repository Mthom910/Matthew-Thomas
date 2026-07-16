#!/usr/bin/env python3
"""Build Revenue Initiative H1 2026 PPTX using Hunter Douglas template."""
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ── DATA ──────────────────────────────────────────────────────────────────────
with open(r'C:\Users\matth\API connection\initiative_data.json') as f:
    d = json.load(f)
k = d['kpis']
buckets = d['buckets']
price_segs = d['price_segs']
top_reps = d['top_reps']
monthly26 = d['monthly26']
monthly25 = d['monthly25']
appts26 = d['monthly_appts26']
appts25 = d['monthly_appts25']

# ── PRESENTATION ──────────────────────────────────────────────────────────────
prs = Presentation(r'C:\Users\matth\OneDrive\work\template.pptx')

# Delete all existing slides
xml_prs = prs.part._element
sldIdLst = xml_prs.find(qn('p:sldIdLst'))
for sldId in list(sldIdLst):
    rId = sldId.get(qn('r:id'))
    sldIdLst.remove(sldId)
    prs.part.drop_rel(rId)

# ── COLORS ────────────────────────────────────────────────────────────────────
SLATE  = RGBColor(0x5A, 0x67, 0x70)   # Hunter Douglas slate gray
TAN    = RGBColor(0xF4, 0xEB, 0xE2)   # Hunter Douglas warm tan
ORANGE = RGBColor(0xF7, 0x96, 0x46)   # Template accent (borders)
TEAL   = RGBColor(0x4B, 0xAC, 0xC6)   # Blue-teal
DARK   = RGBColor(0x3C, 0x3C, 0x3C)   # Near-black
MID    = RGBColor(0x88, 0x88, 0x88)   # Mid gray
GREEN  = RGBColor(0x27, 0xAE, 0x60)   # Positive green
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INIT_BG= RGBColor(0xFF, 0xF3, 0xE6)   # Light orange tint for initiative rows

# ── DIMENSIONS ────────────────────────────────────────────────────────────────
# Slide: 12192000 × 6858000 EMUs = 13.33" × 7.5"
W   = prs.slide_width
H   = prs.slide_height
LM  = Inches(0.22)
CW  = Inches(12.89)
TOP = Inches(1.05)   # Below header
BOT = Inches(6.65)   # Above footer

MONTH_LABELS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']

# ── HELPERS ───────────────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[1])

def set_title(slide, text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            tf = ph.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.italic = True
                    run.font.bold = False
                    run.font.size = Pt(24)
                    run.font.color.rgb = SLATE
            return
    # Fallback
    _txb(slide, text, LM, Inches(0.2), CW, Inches(0.65),
         sz=24, italic=True, color=SLATE)

def _txb(slide, text, left, top, width, height,
         sz=13, bold=False, italic=False, color=DARK,
         align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def _multi_para(slide, lines, left, top, width, height,
                sz=12, color=DARK, bold_first=False, wrap=True,
                line_spacing_pt=None):
    """Add a text box with multiple paragraphs."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return tb

def stat_box(slide, value, label, sub, left, top, w=Inches(2.9),
             val_color=SLATE, val_sz=38):
    """Big number + label + optional sub-label."""
    _txb(slide, value, left, top, w, Inches(0.72),
         sz=val_sz, bold=True, color=val_color, align=PP_ALIGN.CENTER, wrap=False)
    _txb(slide, label, left, top + Inches(0.74), w, Inches(0.38),
         sz=11, bold=True, color=DARK, align=PP_ALIGN.CENTER, wrap=False)
    if sub:
        _txb(slide, sub, left, top + Inches(1.12), w, Inches(0.35),
             sz=10, color=MID, align=PP_ALIGN.CENTER, wrap=False)

def rect(slide, left, top, width, height, fill_color=None,
         line_color=None, line_width_pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_width_pt)
    else:
        shp.line.fill.background()
    return shp

def fmtC(val):
    if val >= 1_000_000: return f'${val/1_000_000:.2f}M'
    if val >= 1_000: return f'${val/1_000:.0f}K'
    return f'${val:.0f}'

def fmtC1(val):
    if val >= 1_000_000: return f'${val/1_000_000:.1f}M'
    if val >= 1_000: return f'${val/1_000:.0f}K'
    return f'${val:.0f}'

def pct_delta(a, b):
    d = a - b
    sign = '+' if d >= 0 else ''
    return f'{sign}{d:.1f}%'

def pp_delta(a, b):
    d = a - b
    sign = '+' if d >= 0 else ''
    return f'{sign}{d:.1f}pp'

def yoy(a, b):
    d = (a - b) / b * 100
    sign = '+' if d >= 0 else ''
    return f'{sign}{d:.1f}%'

def divider(slide, top, color=TAN):
    rect(slide, LM, top, CW, Inches(0.03), fill_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()

# Large title block
rect(sld, LM, Inches(1.4), CW, Inches(2.0), fill_color=SLATE)
_txb(sld, 'Revenue Management Initiative',
     LM + Inches(0.3), Inches(1.55), CW - Inches(0.6), Inches(0.72),
     sz=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT, wrap=False)
_txb(sld, 'H1 2026 Performance Review  ·  January – June 2026',
     LM + Inches(0.3), Inches(2.3), CW - Inches(0.6), Inches(0.5),
     sz=16, color=TAN, align=PP_ALIGN.LEFT, wrap=False)

# Subtitle bar
rect(sld, LM, Inches(3.55), CW, Inches(0.04), fill_color=ORANGE)

# Three callout stats on cover
x_positions = [LM + Inches(0.5), LM + Inches(4.5), LM + Inches(8.5)]
cover_stats = [
    (fmtC1(k['rev26']), '$5K+ Revenue', 'H1 2026'),
    (f"{k['init_pct26']:.0f}%", 'Initiative Adoption', f"35–40% disc bracket"),
    (f"${k['comm26']/1000:.0f}K", 'Total Commission', f"8% base · 6% initiative"),
]
for x, (v, l, s) in zip(x_positions, cover_stats):
    stat_box(sld, v, l, s, x, Inches(3.75), w=Inches(3.5), val_color=SLATE)

# Footer note
_txb(sld, 'Hunter Douglas Australia  ·  Insyte Data  ·  Excludes GST  ·  Service jobs excluded',
     LM, Inches(5.85), CW, Inches(0.4), sz=10, color=MID,
     align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Executive Summary · H1 2026')

# 6 KPI boxes in 2 rows of 3
box_w = Inches(3.9)
box_gap = Inches(0.3)
row1_y = TOP + Inches(0.05)
row2_y = TOP + Inches(1.75)

kpis = [
    (fmtC1(k['rev26']),  f"$5K+ Revenue", f"vs {fmtC1(k['rev25'])} · {yoy(k['rev26'],k['rev25'])} YoY"),
    (str(k['jobs26']),    '$5K+ Orders',   f"vs {k['jobs25']} · {yoy(k['jobs26'],k['jobs25'])} YoY"),
    (fmtC1(k['comm26']), 'Total Commission', f"vs {fmtC1(k['comm25'])} · {yoy(k['comm26'],k['comm25'])} YoY"),
    (f"{k['acts26']:,}",  'Sales Appointments', f"vs {k['acts25']:,} · {yoy(k['acts26'],k['acts25'])} YoY"),
    (f"{k['conv_over26']:.1f}%", '$5K+ Conversion Rate', f"vs {k['conv_over25']:.1f}% · {pp_delta(k['conv_over26'],k['conv_over25'])} YoY"),
    (f"{k['init_pct26']:.0f}%", 'Initiative Adoption', f"35–40% bracket · vs {k['init_pct25']:.0f}% prior year"),
]

for i, (v, l, s) in enumerate(kpis):
    col = i % 3
    row = i // 3
    x = LM + col * (box_w + box_gap)
    y = row1_y if row == 0 else row2_y
    stat_box(sld, v, l, s, x, y, w=box_w, val_color=SLATE)

divider(sld, row2_y + Inches(1.55))

# Initiative achievement line
_txb(sld, f"Initiative: {k['init_jobs26']} jobs at 35–40% discount bracket · "
         f"{fmtC1(buckets[4]['rev26'])} revenue · "
         f"{k['init_pct26']:.0f}% of all $5K+ orders (was {k['init_pct25']:.0f}%)",
     LM, row2_y + Inches(1.65), CW, Inches(0.45),
     sz=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TOTAL ORDER VOLUME
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Total Order Volume · All Price Points')

# Headline KPIs
stat_box(sld, f"{k['all_jobs26']:,}", 'Total Orders (2026)',
         f"vs {k['all_jobs25']:,} · {yoy(k['all_jobs26'],k['all_jobs25'])} YoY",
         LM, TOP, w=Inches(3.8))
stat_box(sld, fmtC1(k['all_rev26']), 'Total Revenue (ex-GST)',
         f"vs {fmtC1(k['all_rev25'])} · {yoy(k['all_rev26'],k['all_rev25'])} YoY",
         LM + Inches(4.2), TOP, w=Inches(3.8))
stat_box(sld, f"{k['aov_all26']:,}", 'Average Order Value',
         f"vs ${k['aov_all25']:,} · {yoy(k['aov_all26'],k['aov_all25'])} YoY",
         LM + Inches(8.4), TOP, w=Inches(3.8))

divider(sld, TOP + Inches(1.55))

# Monthly table
row_h = Inches(0.48)
tbl_top = TOP + Inches(1.7)
col_widths = [Inches(2.3), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['Month', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
hdr_y = tbl_top
rect(sld, LM, hdr_y, CW, row_h, fill_color=SLATE)
for i, (h, x, cw) in enumerate(zip(headers, col_starts, col_widths)):
    _txb(sld, h, x + Inches(0.08), hdr_y + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows_data = [
    ('2026 Orders', [str(m['jobs']) for m in monthly26]),
    ('2026 Revenue', [fmtC1(m['rev']) for m in monthly26]),
    ('2025 Orders', [str(m['jobs']) for m in monthly25]),
    ('2025 Revenue', [fmtC1(m['rev']) for m in monthly25]),
]
for ri, (lbl, vals) in enumerate(rows_data):
    ry = hdr_y + row_h * (ri + 1)
    bg = TAN if ri % 2 == 0 else WHITE
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    _txb(sld, lbl, col_starts[0] + Inches(0.08), ry + Inches(0.1),
         col_widths[0] - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=(ri < 2), color=DARK, align=PP_ALIGN.LEFT)
    for ci, (v, x, cw) in enumerate(zip(vals, col_starts[1:], col_widths[1:])):
        _txb(sld, v, x + Inches(0.08), ry + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
             sz=11, color=DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — $5K+ INITIATIVE ORDERS (MONTHLY)
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, '$5K+ Initiative Orders · Monthly Trend')

# KPIs
stat_box(sld, str(k['jobs26']), '$5K+ Orders H1 2026',
         f"vs {k['jobs25']} · {yoy(k['jobs26'],k['jobs25'])} YoY",
         LM, TOP, w=Inches(3.8))
stat_box(sld, fmtC1(k['rev26']), '$5K+ Revenue H1 2026',
         f"vs {fmtC1(k['rev25'])} · {yoy(k['rev26'],k['rev25'])} YoY",
         LM + Inches(4.2), TOP, w=Inches(3.8))
stat_box(sld, f"{k['aov_over26']:,}", 'Avg Order Value ($5K+)',
         f"vs ${k['aov_over25']:,} · {yoy(k['aov_over26'],k['aov_over25'])} YoY",
         LM + Inches(8.4), TOP, w=Inches(3.8))

divider(sld, TOP + Inches(1.55))

# Monthly table - $5K+ orders and initiative count
row_h = Inches(0.48)
tbl_top = TOP + Inches(1.7)
col_widths = [Inches(2.3), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65), Inches(1.65)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
hdr_y = tbl_top
rect(sld, LM, hdr_y, CW, row_h, fill_color=SLATE)
for i, (h, x, cw) in enumerate(zip(headers, col_starts, col_widths)):
    _txb(sld, h, x + Inches(0.08), hdr_y + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows_data = [
    ('2026 $5K+ Orders', [str(m['jobs']) for m in monthly26]),
    ('2026 $5K+ Revenue', [fmtC1(m['rev']) for m in monthly26]),
    ('Initiative Orders (35–40%)', [str(m['init']) for m in monthly26]),
    ('2025 $5K+ Orders', [str(m['jobs']) for m in monthly25]),
    ('2025 $5K+ Revenue', [fmtC1(m['rev']) for m in monthly25]),
]
for ri, (lbl, vals) in enumerate(rows_data):
    ry = hdr_y + row_h * (ri + 1)
    bg = INIT_BG if ri == 2 else (TAN if ri % 2 == 0 else WHITE)
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    _txb(sld, lbl, col_starts[0] + Inches(0.08), ry + Inches(0.1),
         col_widths[0] - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=(ri == 2), color=ORANGE if ri == 2 else DARK, align=PP_ALIGN.LEFT)
    for ci, (v, x, cw) in enumerate(zip(vals, col_starts[1:], col_widths[1:])):
        _txb(sld, v, x + Inches(0.08), ry + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
             sz=11, bold=(ri == 2), color=ORANGE if ri == 2 else DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — INITIATIVE ADOPTION · DISCOUNT BUCKETS
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Initiative Adoption · Discount Bracket Analysis')

# Headline
_txb(sld, f"35–40% discount bracket: {k['init_jobs26']} orders in H1 2026 "
         f"({k['init_pct26']:.0f}% of all $5K+ orders) vs {k['init_jobs25']} ({k['init_pct25']:.0f}%) in H1 2025",
     LM, TOP, CW, Inches(0.45), sz=13, bold=True, color=DARK)

row_h = Inches(0.55)
tbl_top = TOP + Inches(0.55)
# Columns: Bracket | Initiative? | 2026 Orders | 2026 Revenue | 2026 GP% | 2025 Orders | 2025 Revenue | YoY Orders
col_widths = [Inches(1.8), Inches(1.1), Inches(1.65), Inches(1.65), Inches(1.0), Inches(1.65), Inches(1.65), Inches(1.3)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['Discount', 'Initiative', '2026 Orders', '2026 Revenue', '2026 GP%', '2025 Orders', '2025 Revenue', 'YoY Orders']
rect(sld, LM, tbl_top, CW, row_h, fill_color=SLATE)
for h, x, cw in zip(headers, col_starts, col_widths):
    _txb(sld, h, x + Inches(0.06), tbl_top + Inches(0.1), cw - Inches(0.12), row_h - Inches(0.1),
         sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=False)

for ri, b in enumerate(buckets):
    ry = tbl_top + row_h * (ri + 1)
    is_init = b['initiative']
    bg = INIT_BG if is_init else (TAN if ri % 2 == 0 else WHITE)
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    gp_pct = b['gp26'] / b['rev26'] * 100 if b['rev26'] > 0 else 0
    yoy_orders = f"+{b['count26'] - b['count25']}" if b['count26'] >= b['count25'] else str(b['count26'] - b['count25'])
    txt_color = ORANGE if is_init else DARK
    row_vals = [
        b['label'],
        '★ Yes' if is_init else '—',
        str(b['count26']),
        fmtC1(b['rev26']),
        f"{gp_pct:.0f}%",
        str(b['count25']),
        fmtC1(b['rev25']),
        yoy_orders,
    ]
    for v, x, cw in zip(row_vals, col_starts, col_widths):
        _txb(sld, v, x + Inches(0.06), ry + Inches(0.11), cw - Inches(0.12), row_h - Inches(0.12),
             sz=11, bold=is_init, color=txt_color, align=PP_ALIGN.CENTER, wrap=False)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — $5K+ OPPORTUNITY CONVERSION
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, '$5K+ Opportunity Conversion · Appointment to Order')

# Funnel KPIs
bw = Inches(3.6)
stat_box(sld, f"{k['acts26']:,}", 'Sales Appointments 2026',
         f"vs {k['acts25']:,} · {yoy(k['acts26'],k['acts25'])} YoY",
         LM, TOP, w=bw)
_txb(sld, '→', LM + Inches(3.7), TOP + Inches(0.15), Inches(0.8), Inches(0.7),
     sz=32, color=MID, align=PP_ALIGN.CENTER)
stat_box(sld, str(k['jobs26']), '$5K+ Orders 2026',
         f"vs {k['jobs25']} · {yoy(k['jobs26'],k['jobs25'])} YoY",
         LM + Inches(4.6), TOP, w=bw)
_txb(sld, '=', LM + Inches(8.3), TOP + Inches(0.15), Inches(0.8), Inches(0.7),
     sz=32, color=MID, align=PP_ALIGN.CENTER)
stat_box(sld, f"{k['conv_over26']:.1f}%", 'Conversion Rate 2026',
         f"vs {k['conv_over25']:.1f}% · {pp_delta(k['conv_over26'],k['conv_over25'])} YoY",
         LM + Inches(9.2), TOP, w=bw)

divider(sld, TOP + Inches(1.55))

# Monthly appointments + conversion table
row_h = Inches(0.48)
tbl_top = TOP + Inches(1.7)
col_widths = [Inches(2.5), Inches(1.55), Inches(1.55), Inches(1.55), Inches(1.55), Inches(1.55), Inches(1.55)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

mths = ['2026-01','2026-02','2026-03','2026-04','2026-05','2026-06']
mths25 = ['2025-01','2025-02','2025-03','2025-04','2025-05','2025-06']

rect(sld, LM, tbl_top, CW, row_h, fill_color=SLATE)
for h, x, cw in zip(['', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun'], col_starts, col_widths):
    _txb(sld, h, x + Inches(0.08), tbl_top + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if h else PP_ALIGN.LEFT)

conv26 = [f"{monthly26[i]['jobs'] / appts26[mths[i]] * 100:.1f}%" for i in range(6)]
conv25 = [f"{monthly25[i]['jobs'] / appts25[mths25[i]] * 100:.1f}%" for i in range(6)]

rows_data = [
    ('2026 Appointments', [str(appts26[m]) for m in mths]),
    ('2026 $5K+ Orders', [str(m['jobs']) for m in monthly26]),
    ('2026 Conversion %', conv26),
    ('2025 Appointments', [str(appts25[m]) for m in mths25]),
    ('2025 $5K+ Orders', [str(m['jobs']) for m in monthly25]),
    ('2025 Conversion %', conv25),
]
for ri, (lbl, vals) in enumerate(rows_data):
    ry = tbl_top + row_h * (ri + 1)
    bg = TAN if ri < 3 else (WHITE if ri == 3 else (TAN if ri == 4 else WHITE))
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    is_conv = 'Conversion' in lbl
    _txb(sld, lbl, col_starts[0] + Inches(0.08), ry + Inches(0.1),
         col_widths[0] - Inches(0.16), row_h - Inches(0.1),
         sz=11, bold=is_conv, color=TEAL if is_conv else DARK, align=PP_ALIGN.LEFT)
    for v, x, cw in zip(vals, col_starts[1:], col_widths[1:]):
        _txb(sld, v, x + Inches(0.08), ry + Inches(0.1), cw - Inches(0.16), row_h - Inches(0.1),
             sz=11, bold=is_conv, color=TEAL if is_conv else DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AOV BY PRICE SEGMENT
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Average Order Value · All Orders by Price Segment')

_txb(sld, f"All {k['all_jobs26']:,} orders  ·  Total revenue {fmtC1(k['all_rev26'])}  ·  "
         f"Overall AOV ${k['aov_all26']:,} (vs ${k['aov_all25']:,} H1 2025)",
     LM, TOP, CW, Inches(0.4), sz=12, color=DARK)

row_h = Inches(0.6)
tbl_top = TOP + Inches(0.52)
col_widths = [Inches(2.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.5)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['Segment', '2026 Orders', '2026 Revenue', '2026 AOV', '2025 Orders', '2025 Revenue', '2025 AOV']
rect(sld, LM, tbl_top, CW, row_h, fill_color=SLATE)
for h, x, cw in zip(headers, col_starts, col_widths):
    _txb(sld, h, x + Inches(0.08), tbl_top + Inches(0.12), cw - Inches(0.16), row_h - Inches(0.12),
         sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER if h != 'Segment' else PP_ALIGN.LEFT)

for ri, ps in enumerate(price_segs):
    ry = tbl_top + row_h * (ri + 1)
    bg = TAN if ri % 2 == 0 else WHITE
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    row_vals = [
        ps['label'],
        str(ps['count26']),
        fmtC1(ps['rev26']),
        f"${ps['aov26']:,}",
        str(ps['count25']),
        fmtC1(ps['rev25']),
        f"${ps['aov25']:,}",
    ]
    for v, x, cw in zip(row_vals, col_starts, col_widths):
        _txb(sld, v, x + Inches(0.08), ry + Inches(0.12), cw - Inches(0.16), row_h - Inches(0.12),
             sz=12, color=DARK, align=PP_ALIGN.LEFT if v == ps['label'] else PP_ALIGN.CENTER)

# Total row
ry = tbl_top + row_h * (len(price_segs) + 1)
rect(sld, LM, ry, CW, row_h, fill_color=SLATE)
total_vals = ['Total', str(k['all_jobs26']), fmtC1(k['all_rev26']), f"${k['aov_all26']:,}",
              str(k['all_jobs25']), fmtC1(k['all_rev25']), f"${k['aov_all25']:,}"]
for v, x, cw in zip(total_vals, col_starts, col_widths):
    _txb(sld, v, x + Inches(0.08), ry + Inches(0.12), cw - Inches(0.16), row_h - Inches(0.12),
         sz=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT if v == 'Total' else PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MARGIN IMPACT
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Gross Profit Impact · $5K+ Orders by Discount Bracket')

# GP KPIs
stat_box(sld, fmtC(k['gp26']), 'Total GP 2026',
         f"{k['gp_pct26']:.1f}% GP margin",
         LM, TOP, w=Inches(3.6))
stat_box(sld, fmtC(k['gp25']), 'Total GP 2025',
         f"{k['gp_pct25']:.1f}% GP margin",
         LM + Inches(3.9), TOP, w=Inches(3.6))
stat_box(sld, f"{pp_delta(k['gp_pct26'], k['gp_pct25'])}", 'Margin Change YoY',
         f"Discount pressure from initiative",
         LM + Inches(7.8), TOP, w=Inches(4.4),
         val_color=MID if k['gp_pct26'] < k['gp_pct25'] else GREEN)

divider(sld, TOP + Inches(1.55))

# Table: bucket GP breakdown
row_h = Inches(0.52)
tbl_top = TOP + Inches(1.7)
col_widths = [Inches(1.8), Inches(1.65), Inches(1.65), Inches(1.0), Inches(1.65), Inches(1.65), Inches(1.0), Inches(1.8)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['Bracket', '2026 Revenue', '2026 GP', '2026 GP%', '2025 Revenue', '2025 GP', '2025 GP%', 'GP% Δ']
rect(sld, LM, tbl_top, CW, row_h, fill_color=SLATE)
for h, x, cw in zip(headers, col_starts, col_widths):
    _txb(sld, h, x + Inches(0.05), tbl_top + Inches(0.1), cw - Inches(0.1), row_h - Inches(0.1),
         sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=False)

for ri, b in enumerate(buckets):
    ry = tbl_top + row_h * (ri + 1)
    is_init = b['initiative']
    bg = INIT_BG if is_init else (TAN if ri % 2 == 0 else WHITE)
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    gp26 = b['gp26'] / b['rev26'] * 100 if b['rev26'] > 0 else 0
    gp25 = b['gp25'] / b['rev25'] * 100 if b['rev25'] > 0 else 0
    delta = gp26 - gp25
    txt_color = ORANGE if is_init else DARK
    row_vals = [
        b['label'],
        fmtC1(b['rev26']), fmtC(b['gp26']), f"{gp26:.0f}%",
        fmtC1(b['rev25']), fmtC(b['gp25']), f"{gp25:.0f}%",
        f"{'+' if delta>=0 else ''}{delta:.1f}pp",
    ]
    for v, x, cw in zip(row_vals, col_starts, col_widths):
        _txb(sld, v, x + Inches(0.05), ry + Inches(0.1), cw - Inches(0.1), row_h - Inches(0.1),
             sz=11, bold=is_init, color=txt_color, align=PP_ALIGN.CENTER, wrap=False)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — COMMISSION EARNINGS
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Commission Earnings · H1 2026 Top Consultants')

_txb(sld,
     f"8% base rate on all $5K+ orders  ·  6% on 35–40% bracket (initiative)  ·  "
     f"Total commission pool: {fmtC1(k['comm26'])}  ·  Effective rate: {k['comm_rate26']:.1f}%",
     LM, TOP, CW, Inches(0.45), sz=11, color=DARK)

row_h = Inches(0.52)
tbl_top = TOP + Inches(0.58)
col_widths = [Inches(3.3), Inches(1.4), Inches(1.65), Inches(1.65), Inches(1.4), Inches(1.65), Inches(1.65), Inches(1.0)]
col_starts = [LM]
for cw in col_widths[:-1]:
    col_starts.append(col_starts[-1] + cw)

headers = ['Consultant', '2026 Orders', '2026 Revenue', '2026 GP', 'Init. Orders', '2026 Comm.', '2025 Comm.', 'Δ YoY']
rect(sld, LM, tbl_top, CW, row_h, fill_color=SLATE)
for h, x, cw in zip(headers, col_starts, col_widths):
    _txb(sld, h, x + Inches(0.05), tbl_top + Inches(0.1), cw - Inches(0.1), row_h - Inches(0.1),
         sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER if h != 'Consultant' else PP_ALIGN.LEFT, wrap=False)

for ri, rep in enumerate(top_reps):
    ry = tbl_top + row_h * (ri + 1)
    bg = TAN if ri % 2 == 0 else WHITE
    rect(sld, LM, ry, CW, row_h, fill_color=bg)
    comm_delta = ''
    if rep.get('comm25'):
        d = rep['comm26'] - rep['comm25']
        comm_delta = f"{'+' if d >= 0 else ''}{d/1000:.1f}K"
    row_vals = [
        rep['rep'],
        str(rep['jobs26']),
        fmtC1(rep['rev26']),
        fmtC1(rep['gp26']),
        str(rep.get('init26', '—')),
        fmtC(rep['comm26']),
        fmtC(rep.get('comm25', 0)) if rep.get('comm25') else '—',
        comm_delta or '—',
    ]
    for v, x, cw in zip(row_vals, col_starts, col_widths):
        _txb(sld, v, x + Inches(0.05), ry + Inches(0.1), cw - Inches(0.1), row_h - Inches(0.1),
             sz=11, color=DARK,
             align=PP_ALIGN.LEFT if v == rep['rep'] else PP_ALIGN.CENTER, wrap=False)

# Total row
if len(top_reps) < 8:
    ry = tbl_top + row_h * (len(top_reps) + 1)
    rect(sld, LM, ry, CW, row_h, fill_color=TAN)
    _txb(sld, f'Total commission pool (all consultants): {fmtC1(k["comm26"])} '
              f'vs {fmtC1(k["comm25"])} H1 2025 ({yoy(k["comm26"],k["comm25"])} YoY)',
         LM + Inches(0.1), ry + Inches(0.1), CW - Inches(0.2), row_h - Inches(0.1),
         sz=11, bold=True, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
sld = new_slide()
set_title(sld, 'Key Findings & Recommended Actions')

col_w = Inches(6.1)
col_gap = Inches(0.5)
col2_x = LM + col_w + col_gap

findings = [
    ('Revenue Growth', f"{yoy(k['rev26'],k['rev25'])} growth in $5K+ revenue ({fmtC1(k['rev26'])} vs {fmtC1(k['rev25'])})"),
    ('Initiative Adoption', f"40% of $5K+ orders in 35–40% bracket (up from 31% H1 2025)"),
    ('Order Volume', f"601 orders $5K+ (+9.9% YoY) from 8,030 sales appointments"),
    ('Conversion Uplift', f"7.5% $5K+ conversion rate vs 7.1% prior year (+0.4pp)"),
    ('Margin Pressure', f"GP margin 47.6% vs 49.1% H1 2025 (-1.5pp); discount bracket shift noted"),
    ('Commission Pool', f"{fmtC1(k['comm26'])} earned H1 2026 ({yoy(k['comm26'],k['comm25'])} YoY increase)"),
]

actions = [
    ('Review Q3 Bracket Thresholds', 'Assess whether 35–40% bracket should widen or narrow based on margin data'),
    ('Target Sub-35% Bracket Growth', 'Opportunity to lift GP margin by incentivising 30–35% discount orders'),
    ('High-Value Segment Focus', '$10K+ segment: 146 orders, AOV $14,524 — consultants need enablement support'),
    ('Conversion Training', '7.5% conversion from appointment to $5K+ order — target 8%+ for H2 2026'),
]

rect(sld, LM, TOP, col_w, Inches(0.38), fill_color=SLATE)
_txb(sld, 'Key Findings', LM + Inches(0.1), TOP + Inches(0.05), col_w - Inches(0.2), Inches(0.3),
     sz=13, bold=True, color=WHITE)

y = TOP + Inches(0.5)
for title, detail in findings:
    rect(sld, LM, y, col_w, Inches(0.04), fill_color=TAN)
    _txb(sld, title, LM + Inches(0.08), y + Inches(0.08), col_w - Inches(0.16), Inches(0.3),
         sz=11, bold=True, color=DARK)
    _txb(sld, detail, LM + Inches(0.08), y + Inches(0.38), col_w - Inches(0.16), Inches(0.38),
         sz=11, color=MID, wrap=True)
    y += Inches(0.82)

rect(sld, col2_x, TOP, col_w, Inches(0.38), fill_color=ORANGE)
_txb(sld, 'Recommended Actions', col2_x + Inches(0.1), TOP + Inches(0.05), col_w - Inches(0.2), Inches(0.3),
     sz=13, bold=True, color=WHITE)

y2 = TOP + Inches(0.5)
for title, detail in actions:
    rect(sld, col2_x, y2, col_w, Inches(0.04), fill_color=INIT_BG)
    _txb(sld, title, col2_x + Inches(0.08), y2 + Inches(0.08), col_w - Inches(0.16), Inches(0.3),
         sz=11, bold=True, color=ORANGE)
    _txb(sld, detail, col2_x + Inches(0.08), y2 + Inches(0.38), col_w - Inches(0.16), Inches(0.45),
         sz=11, color=DARK, wrap=True)
    y2 += Inches(1.05)

# ── SAVE ──────────────────────────────────────────────────────────────────────
out = r'C:\Users\matth\API connection\Revenue_Initiative_H1_2026_v5.pptx'
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
