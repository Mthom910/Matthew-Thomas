"""Build Revenue Management Initiative PPTX from initiative_data.json."""
import json, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData, CategoryChartData
from pptx.oxml.ns import qn

with open("initiative_data.json") as f:
    D = json.load(f)
K = D["kpis"]

NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4)
TEAL   = RGBColor(0x02, 0x80, 0x90)
GREEN  = RGBColor(0x05, 0x96, 0x69)
RED    = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x64, 0x74, 0x8B)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)

def fm(v):
    a = abs(v)
    if a >= 1e6: return f"${v/1e6:.2f}M"
    if a >= 1e3: return f"${v/1e3:.0f}K"
    return f"${v:.0f}"
def fp(v): return f"{v:.1f}%"
def fn(v): return f"{v:,.0f}"
def yoy(a, b): d = (a-b)/b*100; return f"{'+'if d>=0 else ''}{d:.1f}% YoY"
def pp_delta(a, b): d = a-b; return f"{'+'if d>=0 else ''}{d:.1f}pp"
def dollar_delta(a, b): d = a-b; sign="+"; sign = "+" if d>=0 else "-"; return f"{sign}{fm(abs(d))}"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None, lw=0):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else: sh.line.fill.background()
    return sh

def add_tb(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
           italic=False, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap; tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def add_header(slide, title, subtitle, bg=NAVY):
    add_rect(slide, 0, 0, 13.33, 0.85, bg)
    add_tb(slide, 0.4, 0.07, 10, 0.5, title, 24, WHITE, bold=True, font="Cambria")
    add_tb(slide, 0.4, 0.55, 12, 0.28, subtitle, 9, RGBColor(0xA0,0xC4,0xE8))

def add_kpi(slide, x, y, w, h, label, value, sub, delta, up, vcol, bg=None, border=None):
    if bg is None: bg = LIGHT
    if border is None: border = RGBColor(0xE2,0xE8,0xF0)
    add_rect(slide, x, y, w, h, bg, border, 0.5)
    add_tb(slide, x+.08, y+.08, w-.16, .24, label, 8.5, SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(slide, x+.08, y+.32, w-.16, .55, value, 28, vcol, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    add_tb(slide, x+.08, y+.86, w-.16, .20, sub, 7.5, SLATE, align=PP_ALIGN.CENTER)
    add_tb(slide, x+.08, y+1.06, w-.16, .20, delta, 8.5, GREEN if up else RED, bold=True, align=PP_ALIGN.CENTER)

def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def add_chart(slide, x, y, w, h, series, cats, colors=None,
              ctype=XL_CHART_TYPE.COLUMN_CLUSTERED, legend=True, vals=False):
    cd = CategoryChartData(); cd.categories = cats
    for nm, data in series: cd.add_series(nm, data)
    cs = slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = cs.chart; ch.has_title = False; ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    if colors:
        for ser, col in zip(ch.series, colors):
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
    if vals:
        for ser in ch.series:
            ser.data_labels.show_value = True
            ser.data_labels.font.size = Pt(8); ser.data_labels.font.color.rgb = NAVY
    for ax in [ch.category_axis, ch.value_axis]:
        ax.tick_labels.font.size = Pt(8); ax.tick_labels.font.color.rgb = SLATE
    ch.value_axis.has_major_gridlines = True
    ch.category_axis.has_major_gridlines = False
    return ch

# ── Shorthand data ─────────────────────────────────────────
mLabels  = ["Jan","Feb","Mar","Apr","May","Jun"]
mRevs26  = [round(m["rev"]/1000,1) for m in D["monthly26"]]
mRevs25  = [round(m["rev"]/1000,1) for m in D["monthly25"]]
mJobs26  = [m["jobs"] for m in D["monthly26"]]
mJobs25  = [m["jobs"] for m in D["monthly25"]]
mInits26 = [m["init"] for m in D["monthly26"]]
mAppts26 = list(D["monthly_appts26"].values())
mAppts25 = list(D["monthly_appts25"].values())

rev_growth   = round((K["rev26"]-K["rev25"])/K["rev25"]*100,1)
init_growth  = round((K["init_jobs26"]-K["init_jobs25"])/K["init_jobs25"]*100)
gp_growth    = round((K["gp26"]-K["gp25"])/K["gp25"]*100,1)
comm_diff    = K["comm26"]-K["comm25"]
init_b       = D["buckets"][4]
prev_b       = D["buckets"][3]
init_gp_pct  = round(init_b["gp26"]/init_b["rev26"]*100,1)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — COVER
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 0.55, 7.5, TEAL)
circ = s.shapes.add_shape(9, Inches(9.5), Inches(-1.0), Inches(5.0), Inches(5.0))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0x02,0x80,0x90); circ.line.fill.background()

add_rect(s, 0.8, 1.1, 3.0, 0.35, TEAL)
add_tb(s, 0.8, 1.1, 3.0, 0.35, "SENIOR MANAGEMENT BRIEFING", 8, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_tb(s, 0.8, 1.6, 9.5, 1.1, "Revenue Management", 52, WHITE, bold=True, font="Cambria")
add_tb(s, 0.8, 2.65, 9.5, 0.7, "Initiative", 52, RGBColor(0xA0,0xC4,0xE8), bold=True, font="Cambria")
add_tb(s, 0.8, 3.55, 11, 0.45, "YTD Performance Review  ·  January – June 2026 vs January – June 2025", 14,
       RGBColor(0xA0,0xC4,0xE8))
add_tb(s, 0.8, 4.0, 11, 0.35, "35–40% Commission Bracket Initiative  ·  $5K+ Job Analysis  ·  Ex-GST Confirmed Orders", 10.5,
       RGBColor(0x7A,0xA8,0xCC))

kpi_items = [
    ("Total Order Intake",   fm(K["all_rev26"]),   yoy(K["all_rev26"],K["all_rev25"])),
    ("$5K+ Revenue",         fm(K["rev26"]),        yoy(K["rev26"],K["rev25"])),
    ("Initiative Bracket",   fp(K["init_pct26"]),   f"{K['init_jobs26']} jobs (35–40%)"),
    ("Gross Profit",         fm(K["gp26"]),          fp(K["gp_pct26"])+" GP%"),
]
for i,(lbl,val,sub) in enumerate(kpi_items):
    bx = 0.8 + i*3.08
    add_rect(s, bx, 5.05, 2.8, 1.05, RGBColor(0x1E,0x3A,0x5F), RGBColor(0x7A,0xA8,0xCC), 0.5)
    add_tb(s, bx+.1, 5.1, 2.6, .25, lbl, 8, RGBColor(0xA0,0xC4,0xE8), bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, bx+.1, 5.33, 2.6, .42, val, 22, WHITE, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    add_tb(s, bx+.1, 5.75, 2.6, .28, sub, 8, RGBColor(0x7A,0xA8,0xCC), align=PP_ALIGN.CENTER)

add_tb(s, 0.8, 6.85, 11, 0.25,
       "Prepared by Sales Analytics  ·  Data as at 27 June 2026  ·  All figures Ex-GST  ·  Confirmed orders",
       7.5, SLATE, align=PP_ALIGN.CENTER)

add_notes(s,
    f"H1 2026 Revenue Management Initiative — Senior Management Briefing.\n\n"
    f"Total order intake: {fm(K['all_rev26'])} across {fn(K['all_jobs26'])} jobs ({yoy(K['all_rev26'],K['all_rev25'])}).\n"
    f"$5K+ segment: {fm(K['rev26'])} on {K['jobs26']} jobs ({yoy(K['rev26'],K['rev25'])}).\n"
    f"Initiative bracket adoption: {K['init_pct25']}% → {K['init_pct26']}% of $5K+ jobs.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — EXECUTIVE SUMMARY
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Executive Summary",
           "$5K+ Confirmed Orders  ·  Ex-GST  ·  Jan–Jun 2026 vs Jan–Jun 2025")

kpis = [
    ("$5K+ Revenue",        fm(K["rev26"]),      f"vs {fm(K['rev25'])} (2025)", yoy(K["rev26"],K["rev25"]),       K["rev26"]>K["rev25"],  BLUE),
    ("Jobs Won ($5K+)",     fn(K["jobs26"]),      f"vs {fn(K['jobs25'])} (2025)", yoy(K["jobs26"],K["jobs25"]),   K["jobs26"]>K["jobs25"], GREEN),
    ("Avg Discount (Won)",  fp(K["avg_disc26"]),  f"vs {fp(K['avg_disc25'])} (2025)", pp_delta(K["avg_disc26"],K["avg_disc25"])+" YoY", False, AMBER),
    ("Gross Profit %",      fp(K["gp_pct26"]),    f"vs {fp(K['gp_pct25'])} (2025)", pp_delta(K["gp_pct26"],K["gp_pct25"])+" YoY", False, RED),
]
for i,(lbl,val,sub,delta,up,col) in enumerate(kpis):
    add_kpi(s, 0.3+i*3.28, 1.0, 3.0, 1.5, lbl, val, sub, delta, up, col)

add_rect(s, 0.3, 2.7, 12.73, 4.5, RGBColor(0xEF,0xF6,0xFF), BLUE, 1.5)
add_tb(s, 0.55, 2.82, 3, 0.3, "Key Findings", 11, NAVY, bold=True)

findings = [
    f"1.  Total order intake reached {fm(K['all_rev26'])} across {fn(K['all_jobs26'])} jobs ({yoy(K['all_rev26'],K['all_rev25'])}). The $5K+ segment drove {fm(K['rev26'])} ({yoy(K['rev26'],K['rev25'])}).",
    f"2.  Initiative bracket (35–40% discount) adoption surged from {K['init_pct25']}% to {K['init_pct26']}% of won $5K+ jobs — {K['init_jobs26']} jobs vs {K['init_jobs25']} in 2025 (+{init_growth}%).",
    f"3.  $5K+ appointment-to-order conversion rate improved from {K['conv_over25']}% to {K['conv_over26']}% on {fn(K['acts26'])} sales appointments (+{fn(K['acts26']-K['acts25'])} vs 2025).",
    f"4.  GP% eased {abs(K['gp_pct26']-K['gp_pct25']):.1f}pp to {K['gp_pct26']}% — gross profit dollars still grew {gp_growth}% to {fm(K['gp26'])}. Volume offsets the rate compression.",
]
for i,t in enumerate(findings):
    add_tb(s, 0.55, 3.2+i*.72, 12.35, 0.65, t, 10.5, NAVY)

add_notes(s,
    f"Executive summary.\n"
    f"1. Revenue: all jobs {fm(K['all_rev26'])}, $5K+ jobs {fm(K['rev26'])}, both up YoY.\n"
    f"2. Initiative adoption: {K['init_pct25']}% → {K['init_pct26']}% — 2 in 5 $5K+ jobs now use the bracket.\n"
    f"3. Conversion: {K['conv_over25']}% → {K['conv_over26']}% on $5K+ jobs. More appointments, better conversion.\n"
    f"4. GP%: -{abs(K['gp_pct26']-K['gp_pct25']):.1f}pp but GP$ grew {gp_growth}%.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — ORDER INTAKE PERFORMANCE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Order Intake Performance",
           "Monthly revenue trend (ex-GST)  ·  $5K+ confirmed orders  ·  2026 vs 2025")

add_chart(s, 0.3, 0.95, 8.5, 4.5,
    [("2026 Revenue ($K)", mRevs26), ("2025 Revenue ($K)", mRevs25)],
    mLabels, colors=[BLUE, RGBColor(0x94,0xA3,0xB8)])

panel_x = 9.1
stats = [
    ("$5K+ Rev 2026",   fm(K["rev26"]),                           GREEN),
    ("vs 2025",         yoy(K["rev26"],K["rev25"]),                GREEN),
    ("Total Intake",    fm(K["all_rev26"]),                        BLUE),
    ("vs 2025",         yoy(K["all_rev26"],K["all_rev25"]),        BLUE),
    ("Best Month",      f"{mLabels[mRevs26.index(max(mRevs26))]} — {fm(max(mRevs26)*1000)}", NAVY),
]
for i,(lbl,val,col) in enumerate(stats):
    ry = 1.0+i*.88
    add_rect(s, panel_x, ry, 3.93, .78, LIGHT, RGBColor(0xE2,0xE8,0xF0), 0.5)
    add_tb(s, panel_x+.1, ry+.06, 3.73, .22, lbl, 8, SLATE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, panel_x+.1, ry+.28, 3.73, .38, val, 18, col, bold=True, align=PP_ALIGN.CENTER, font="Cambria")

add_rect(s, 0.3, 5.55, 8.5, .32, NAVY)
for ci,(txt,wx) in enumerate([("Month",1.2),("2026 Jobs",2.0),("2025 Jobs",2.0),("Init. Jobs",3.3)]):
    cx = 0.3+sum([1.2,2.0,2.0,3.3][:ci])
    add_tb(s, cx+.05, 5.57, wx-.1, .28, txt, 8, WHITE, bold=True)

cols_x = [0.3,1.5,3.5,5.5]
for ri,(m26,m25) in enumerate(zip(D["monthly26"],D["monthly25"])):
    mo = int(m26["month"].split("-")[1])
    lbl = ["Jan","Feb","Mar","Apr","May","Jun"][mo-1]
    bg = LIGHT if ri%2==0 else WHITE
    ry = 5.87+ri*.24
    add_rect(s, 0.3, ry, 8.5, .24, bg)
    vals = [lbl, str(m26["jobs"]), str(m25["jobs"]), str(m26["init"])]
    clrs = [NAVY, GREEN if m26["jobs"]>=m25["jobs"] else RED, SLATE, PURPLE]
    for ci,(v,col) in enumerate(zip(vals,clrs)):
        add_tb(s, cols_x[ci]+.05, ry+.02, 1.8, .22, v, 8.5, col, bold=(ci>0))

add_notes(s,
    f"Order intake H1 2026.\n"
    f"$5K+ total: {fm(K['rev26'])} vs {fm(K['rev25'])} — {yoy(K['rev26'],K['rev25'])}\n"
    f"All-jobs total: {fm(K['all_rev26'])} vs {fm(K['all_rev25'])} — {yoy(K['all_rev26'],K['all_rev25'])}\n" +
    "\n".join(f"  {mLabels[i]}: {fm(mRevs26[i]*1000)} 2026 vs {fm(mRevs25[i]*1000)} 2025" for i in range(6)))


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — INITIATIVE BRACKET ADOPTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Initiative Bracket Adoption",
           "35–40% discount bracket  ·  $5K+ jobs  ·  Job count comparison 2026 vs 2025",
           bg=PURPLE)

add_rect(s, 0.3, 1.0, 5.8, 3.0, RGBColor(0xF5,0xF3,0xFF), PURPLE, 1.5)
add_tb(s, 0.5, 1.1, 5.4, .3, "2026 — Initiative Jobs", 11, PURPLE, bold=True, align=PP_ALIGN.CENTER)
add_tb(s, 0.5, 1.35, 5.4, 1.3, str(K["init_jobs26"]), 88, PURPLE, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
add_tb(s, 0.5, 2.62, 5.4, .25, "jobs in 35–40% discount bracket", 11, PURPLE, align=PP_ALIGN.CENTER)
add_tb(s, 0.5, 2.9, 5.4, .25,
       f"{fp(K['init_pct26'])} of all $5K+ won jobs  ·  +{init_growth}% vs 2025", 10,
       PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_rect(s, 6.5, 1.0, 5.8, 3.0, LIGHT, RGBColor(0xCB,0xD5,0xE1), 1.0)
add_tb(s, 6.7, 1.1, 5.4, .3, "2025 — Initiative Jobs (Baseline)", 11, SLATE, bold=True, align=PP_ALIGN.CENTER)
add_tb(s, 6.7, 1.35, 5.4, 1.3, str(K["init_jobs25"]), 88, RGBColor(0x94,0xA3,0xB8),
       bold=True, align=PP_ALIGN.CENTER, font="Cambria")
add_tb(s, 6.7, 2.62, 5.4, .25, "jobs in 35–40% discount bracket", 11, SLATE, align=PP_ALIGN.CENTER)
add_tb(s, 6.7, 2.9, 5.4, .25,
       f"{fp(K['init_pct25'])} of all $5K+ won jobs (prior year baseline)", 10, SLATE, align=PP_ALIGN.CENTER)

add_rect(s, 0.3, 4.2, 12.73, .65, RGBColor(0xED,0xE9,0xFE), PURPLE, 0.8)
add_tb(s, 0.5, 4.28, 12.33, .5,
       f"▲  Adoption grew +{init_growth}%  ·  Revenue in bracket: {fm(init_b['rev25'])} → {fm(init_b['rev26'])} (+{round((init_b['rev26']-init_b['rev25'])/init_b['rev25']*100)}%)  ·  Now 2 in 5 won $5K+ jobs",
       11, PURPLE, bold=True, align=PP_ALIGN.CENTER)

add_chart(s, 0.3, 5.0, 12.73, 2.25,
    [("All Other $5K+ Jobs", [m26-mi for m26,mi in zip(mJobs26,mInits26)]),
     ("Initiative Bracket (35–40%)", mInits26)],
    mLabels, colors=[RGBColor(0x93,0xC5,0xFD), PURPLE],
    ctype=XL_CHART_TYPE.COLUMN_STACKED)

add_notes(s,
    f"Initiative bracket adoption.\n"
    f"{K['init_jobs26']} jobs in 2026 vs {K['init_jobs25']} in 2025 — +{init_growth}%.\n"
    f"Share of $5K+ jobs: {K['init_pct25']}% → {K['init_pct26']}% — nearly 2 in 5 high-value orders now use the bracket.\n"
    f"Revenue in bracket: {fm(init_b['rev25'])} → {fm(init_b['rev26'])} (+{round((init_b['rev26']-init_b['rev25'])/init_b['rev25']*100)}%)\n"
    "Monthly chart shows consistent adoption — not a one-off spike.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — DISCOUNT BUCKET DISTRIBUTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Discount Bracket Distribution",
           "Jobs count and revenue shift by discount band  ·  $5K+ jobs  ·  2026 vs 2025")

bLabels  = [b["label"] for b in D["buckets"]]
bCount26 = [b["count26"] for b in D["buckets"]]
bCount25 = [b["count25"] for b in D["buckets"]]
bRev26k  = [round(b["rev26"]/1000,1) for b in D["buckets"]]
bRev25k  = [round(b["rev25"]/1000,1) for b in D["buckets"]]

add_tb(s, 0.3, 0.92, 6.3, .3, "Jobs Won by Bracket", 10.5, NAVY, bold=True)
add_chart(s, 0.3, 1.2, 6.3, 3.6,
    [("2026", bCount26), ("2025", bCount25)], bLabels,
    colors=[PURPLE, RGBColor(0x94,0xA3,0xB8)], vals=True)

add_tb(s, 6.9, 0.92, 6.1, .3, "Revenue by Bracket ($K, ex-GST)", 10.5, NAVY, bold=True)
add_chart(s, 6.9, 1.2, 6.1, 3.6,
    [("2026", bRev26k), ("2025", bRev25k)], bLabels,
    colors=[BLUE, RGBColor(0x94,0xA3,0xB8)])

prev_change = round((prev_b["count26"]-prev_b["count25"])/prev_b["count25"]*100)
add_rect(s, 0.3, 5.0, 12.73, 2.25, RGBColor(0xF5,0xF3,0xFF), PURPLE, 1.0)
add_tb(s, 0.5, 5.1, 12.33, .3, "★  35–40% Initiative Bracket Spotlight", 11, PURPLE, bold=True)
spotlight = (f"Jobs: {K['init_jobs25']} → {K['init_jobs26']} (+{init_growth}%)   ·   "
             f"Revenue: {fm(init_b['rev25'])} → {fm(init_b['rev26'])} (+{round((init_b['rev26']-init_b['rev25'])/init_b['rev25']*100)}%)   ·   "
             f"GP%: {round(init_b['gp25']/init_b['rev25']*100,1)}% → {init_gp_pct}%   ·   "
             f"Avg Job Value: {fm(init_b['rev26']/init_b['count26'])}")
add_tb(s, 0.5, 5.45, 12.33, .42, spotlight, 10, NAVY, bold=True)
add_tb(s, 0.5, 5.92, 12.33, .58,
    f"The adjacent 30–35% bracket contracted from {prev_b['count25']} to {prev_b['count26']} jobs ({prev_change}%), "
    "confirming consultants shifted upward into the initiative bracket — the intended behaviour driven by enhanced commission.",
    9.5, RGBColor(0x4B,0x55,0x63), italic=True)

add_notes(s,
    f"Discount bracket distribution.\n"
    f"30–35% bracket: {prev_b['count25']} → {prev_b['count26']} jobs ({prev_change}%).\n"
    f"35–40% bracket: {K['init_jobs25']} → {K['init_jobs26']} jobs (+{init_growth}%).\n"
    "This shift is the intended mechanism — reps are using the enhanced commission to offer competitive discounts and close.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — CONVERSION ANALYSIS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Appointment-to-Order Conversion",
           "Sales appointments only (DC calendar bookings)  ·  Won jobs as %  ·  Jan–Jun comparison",
           bg=TEAL)

conv_items = [
    ("All Orders", K["conv_all26"], K["conv_all25"], TEAL),
    ("$5K+ Orders", K["conv_over26"], K["conv_over25"], BLUE),
    ("<$5K Orders", K["conv_under26"], K["conv_under25"], AMBER),
]
for i,(lbl,v26,v25,col) in enumerate(conv_items):
    cx = 0.3+i*4.28
    add_rect(s, cx, 1.0, 4.0, 3.0, LIGHT, RGBColor(0xE2,0xE8,0xF0), 0.8)
    add_tb(s, cx+.15, 1.1, 3.7, .3, lbl, 11, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, cx+.15, 1.4, 3.7, .9, fp(v26), 58, col, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    add_tb(s, cx+.15, 2.32, 3.7, .25, "2026 conversion rate", 8.5, SLATE, align=PP_ALIGN.CENTER)
    add_tb(s, cx+.15, 2.58, 3.7, .25, f"vs {fp(v25)} in 2025", 9, SLATE, align=PP_ALIGN.CENTER)
    d = v26-v25
    add_tb(s, cx+.15, 2.85, 3.7, .32, ("▲ " if d>=0 else "▼ ")+fp(abs(d)),
           12, GREEN if d>=0 else RED, bold=True, align=PP_ALIGN.CENTER)

add_rect(s, 0.3, 4.2, 5.8, 2.1, LIGHT, RGBColor(0xE2,0xE8,0xF0), 0.8)
add_tb(s, 0.5, 4.32, 5.4, .3, "Sales Appointments (DC Calendar)", 11, NAVY, bold=True)
add_tb(s, 0.5, 4.7,  5.4, .35, f"2026:  {fn(K['acts26'])} appointments", 11, NAVY)
add_tb(s, 0.5, 5.08, 5.4, .35, f"2025:  {fn(K['acts25'])} appointments", 11, SLATE)
act_delta = K['acts26']-K['acts25']
add_tb(s, 0.5, 5.46, 5.4, .55,
       f"Change:  +{fn(act_delta)} appointments ({yoy(K['acts26'],K['acts25'])})",
       11, GREEN, bold=True)

add_rect(s, 6.5, 4.2, 6.53, 2.1, RGBColor(0xEC,0xFD,0xF5), GREEN, 1.0)
add_tb(s, 6.7, 4.32, 6.13, .3, "Conversion Quality Insight", 11, GREEN, bold=True)
insight = (
    f"Both appointment volume and conversion rate improved in 2026. "
    f"With {fn(K['acts26'])} sales appointments (+{fn(act_delta)} vs 2025), the team converted "
    f"{K['jobs26']} $5K+ orders at a {K['conv_over26']}% rate — up from {K['conv_over25']}% in 2025. "
    "The initiative's competitive discount position is enabling consultants to close higher-value customers more effectively."
)
add_tb(s, 6.7, 4.7, 6.13, 1.5, insight, 10, RGBColor(0x1F,0x40,0x33))

# Monthly conversion chart
mConv5K26 = [round(j/a*100,1) if a else 0 for j,a in zip(mJobs26,mAppts26)]
mConv5K25 = [round(j/a*100,1) if a else 0 for j,a in zip(mJobs25,mAppts25)]

add_chart(s, 0.3, 6.45, 12.73, 0.85,
    [("2026 $5K+ Conv%", mConv5K26), ("2025 $5K+ Conv%", mConv5K25)],
    mLabels, colors=[BLUE, RGBColor(0x94,0xA3,0xB8)],
    ctype=XL_CHART_TYPE.LINE, legend=True)

add_notes(s,
    f"Conversion analysis — CORRECTED to use sales appointments only (DC calendar bookings).\n\n"
    f"2026: {fn(K['acts26'])} sales appointments (+{fn(act_delta)} vs 2025).\n"
    f"$5K+ conversion: {K['conv_over25']}% → {K['conv_over26']}% — improving despite more appointments.\n"
    f"All-orders conversion: {K['conv_all25']}% → {K['conv_all26']}%.\n\n"
    "Previous data incorrectly included all activity types (calls, emails, notes). "
    "This slide now uses DC sales appointment bookings only.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — AVERAGE ORDER VALUE BY PRICE SEGMENT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Average Order Value by Price Segment",
           "All confirmed orders  ·  Ex-GST  ·  6 price bands  ·  2026 vs 2025")

ps = D["price_segs"]
seg_labels = [p["label"] for p in ps]
aov26 = [p["aov26"] for p in ps]
aov25 = [p["aov25"] for p in ps]
cnt26 = [p["count26"] for p in ps]
cnt25 = [p["count25"] for p in ps]
rev26s = [p["rev26"] for p in ps]

add_tb(s, 0.3, 0.92, 8.5, .3, "Average Order Value — 2026 vs 2025 (ex-GST)", 10.5, NAVY, bold=True)
add_chart(s, 0.3, 1.2, 8.5, 3.5,
    [("2026 AOV", aov26), ("2025 AOV", aov25)],
    seg_labels, colors=[BLUE, RGBColor(0x94,0xA3,0xB8)])

# Table
add_rect(s, 0.3, 4.85, 12.73, .35, NAVY)
col_hdrs = ["Segment","2026 Jobs","2026 AOV","2025 Jobs","2025 AOV","AOV Change","2026 Revenue"]
col_ws   = [1.5,1.2,1.4,1.2,1.4,1.4,2.63]
cx = 0.3
for h,w in zip(col_hdrs,col_ws):
    add_tb(s, cx+.05, 4.87, w-.1, .3, h, 8, WHITE, bold=True); cx+=w

row_h = .33
for ri,p in enumerate(ps):
    ry = 5.2+ri*row_h
    bg = LIGHT if ri%2==0 else WHITE
    add_rect(s, 0.3, ry, 12.73, row_h, bg)
    d = p["aov26"]-p["aov25"]
    vals = [
        (p["label"],          NAVY,  True),
        (fn(p["count26"]),    GREEN if p["count26"]>=p["count25"] else RED, False),
        (fm(p["aov26"]),      BLUE,  True),
        (fn(p["count25"]),    SLATE, False),
        (fm(p["aov25"]),      SLATE, False),
        (("▲ " if d>=0 else "▼ ")+fm(abs(d)), GREEN if d>=0 else RED, True),
        (fm(p["rev26"]),      GREEN, False),
    ]
    cx = 0.3
    for (v,col,bld),w in zip(vals,col_ws):
        add_tb(s, cx+.05, ry+.04, w-.1, row_h-.08, v, 8.5, col, bold=bld); cx+=w

add_notes(s,
    f"Average order value by price segment — all confirmed orders.\n\n"
    "Key observations:\n" +
    "\n".join(f"  {p['label']}: {p['count26']} jobs, AOV ${p['aov26']:,} (vs ${p['aov25']:,} in 2025)" for p in ps) +
    f"\n\nHighest-value segment ($10K+): {ps[5]['count26']} jobs vs {ps[5]['count25']} — +{round((ps[5]['count26']-ps[5]['count25'])/ps[5]['count25']*100)}% job growth at {fm(ps[5]['aov26'])} AOV.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 — MARGIN IMPACT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Gross Profit & Margin Impact",
           "Revenue and GP$ by discount bracket  ·  $5K+ jobs  ·  Ex-GST")

gp_cards = [
    ("GP$ 2026",     fm(K["gp26"]),      f"vs {fm(K['gp25'])} 2025",   yoy(K["gp26"],K["gp25"]),         K["gp26"]>K["gp25"],   GREEN),
    ("GP% 2026",     fp(K["gp_pct26"]),  f"vs {fp(K['gp_pct25'])} 2025", pp_delta(K["gp_pct26"],K["gp_pct25"])+" YoY", False,   AMBER),
    ("$5K+ Rev 2026",fm(K["rev26"]),     f"vs {fm(K['rev25'])} 2025",   yoy(K["rev26"],K["rev25"]),        K["rev26"]>K["rev25"], BLUE),
]
for i,(lbl,val,sub,delta,up,col) in enumerate(gp_cards):
    add_kpi(s, 0.3+i*4.35, 1.0, 4.1, 1.45, lbl, val, sub, delta, up, col)

col_hdrs = ["Bracket","2026 Jobs","2026 Revenue","Rev vs '25","2026 GP$","2026 GP%","2025 GP%","GP% Δ"]
col_ws   = [1.7,0.9,1.3,1.1,1.2,0.9,0.9,0.9]
cx = 0.3
add_rect(s, 0.3, 2.6, 12.73, .35, NAVY)
for h,w in zip(col_hdrs,col_ws):
    add_tb(s, cx+.05, 2.63, w-.1, .3, h, 8, WHITE, bold=True); cx+=w

row_h = .38
for ri,b in enumerate(D["buckets"]):
    ry = 2.95+ri*row_h
    gp26p = round(b["gp26"]/b["rev26"]*100,1) if b["rev26"] else 0
    gp25p = round(b["gp25"]/b["rev25"]*100,1) if b["rev25"] else 0
    gpd   = round(gp26p-gp25p,1)
    revd  = b["rev26"]-b["rev25"]
    bg = RGBColor(0xF5,0xF3,0xFF) if b.get("initiative") else (LIGHT if ri%2==0 else WHITE)
    add_rect(s, 0.3, ry, 12.73, row_h, bg)
    row_vals = [
        (b["label"]+(" ★" if b.get("initiative") else ""), PURPLE if b.get("initiative") else NAVY, b.get("initiative",False)),
        (str(b["count26"]), GREEN if b["count26"]>b["count25"] else RED, False),
        (fm(b["rev26"]), GREEN, False),
        (("+" if revd>=0 else "-")+fm(abs(revd)), GREEN if revd>=0 else RED, True),
        (fm(b["gp26"]), NAVY, False),
        (fp(gp26p), NAVY, False),
        (fp(gp25p), SLATE, False),
        (("+" if gpd>=0 else "")+f"{gpd}pp", GREEN if gpd>=0 else RED, True),
    ]
    cx = 0.3
    for (v,col,bld),w in zip(row_vals,col_ws):
        add_tb(s, cx+.05, ry+.04, w-.1, row_h-.08, v, 8.5, col, bold=bld); cx+=w

add_notes(s,
    f"Margin analysis.\n"
    f"GP$: {fm(K['gp25'])} → {fm(K['gp26'])} (+{gp_growth}% YoY).\n"
    f"GP%: {K['gp_pct25']}% → {K['gp_pct26']}% (-{abs(K['gp_pct26']-K['gp_pct25']):.1f}pp).\n"
    f"Initiative bracket GP%: {round(init_b['gp26']/init_b['rev26']*100,1)}% on {fm(init_b['rev26'])} revenue.\n"
    "Volume growth more than offsets the margin rate compression.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 — COMMISSION EARNINGS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_header(s, "Consultant Commission Earnings",
           "8% base rate  ·  6% on 35–40% bracket jobs  ·  $5K+ orders  ·  Ex-GST")

eff_uplift = round(K["comm_rate26"]-8.0,1)
comm_cards = [
    ("Total Commission 2026", fm(K["comm26"]),  f"vs {fm(K['comm25'])} 2025",         yoy(K["comm26"],K["comm25"]), True,  GREEN),
    ("YoY Increase",          fm(comm_diff),    "additional consultant earnings",       f"+{round(comm_diff/K['comm25']*100)}% YoY", True, BLUE),
    ("Effective Blended Rate",fp(K["comm_rate26"]),"blended rate on $5K+ revenue",    f"+{eff_uplift}pp vs 8% base",True, PURPLE),
]
for i,(lbl,val,sub,delta,up,col) in enumerate(comm_cards):
    add_kpi(s, 0.3+i*4.35, 1.0, 4.1, 1.42, lbl, val, sub, delta, up, col)

rep_cols = ["Consultant","Jobs '26","Revenue '26","Commission '26","Commission '25","YoY Change","Init Jobs"]
rep_cw   = [2.6,0.8,1.35,1.45,1.45,1.2,0.88]
cx = 0.3
add_rect(s, 0.3, 2.55, 12.73, .38, NAVY)
for col,cw in zip(rep_cols,rep_cw):
    add_tb(s, cx+.05, 2.58, cw-.1, .32, col, 8, WHITE, bold=True); cx+=w

row_h = .44
for ri,rep in enumerate(D["top_reps"]):
    ry = 2.93+ri*row_h
    bg = LIGHT if ri%2==0 else WHITE
    add_rect(s, 0.3, ry, 12.73, row_h, bg)
    comm25 = rep.get("comm25",0)
    comm_d = rep["comm26"]-comm25
    name   = rep["rep"].split("(")[0].strip()
    vals = [
        (name, NAVY, True),
        (str(rep["jobs26"]), NAVY, False),
        (fm(rep["rev26"]), GREEN, False),
        (fm(rep["comm26"]), BLUE, True),
        (fm(comm25) if comm25 else "—", SLATE, False),
        ((("+" if comm_d>=0 else "-")+fm(abs(comm_d))) if comm25 else "—", GREEN if comm_d>=0 else RED, bool(comm25)),
        (str(rep["init26"]), PURPLE if rep["init26"]>=10 else SLATE, rep["init26"]>=10),
    ]
    cx = 0.3
    for (v,col,bld),cw in zip(vals,rep_cw):
        add_tb(s, cx+.05, ry+.06, cw-.1, row_h-.1, v, 8.5, col, bold=bld); cx+=cw

tot_rev26  = sum(r["rev26"] for r in D["top_reps"])
tot_comm26 = sum(r["comm26"] for r in D["top_reps"])
tot_comm25 = sum(r.get("comm25",0) for r in D["top_reps"])
tot_init26 = sum(r["init26"] for r in D["top_reps"])
tot_jobs26 = sum(r["jobs26"] for r in D["top_reps"])
tot_d = tot_comm26-tot_comm25

ry = 2.93+len(D["top_reps"])*row_h
add_rect(s, 0.3, ry, 12.73, row_h, RGBColor(0xEF,0xF6,0xFF), NAVY, 0.5)
tot_vals = [
    ("Top 8 Total",NAVY,True),(str(tot_jobs26),NAVY,True),(fm(tot_rev26),GREEN,True),
    (fm(tot_comm26),BLUE,True),(fm(tot_comm25),SLATE,True),
    (("+" if tot_d>=0 else "-")+fm(abs(tot_d)),GREEN if tot_d>=0 else RED,True),
    (str(tot_init26),PURPLE,True),
]
cx = 0.3
for (v,col,bld),cw in zip(tot_vals,rep_cw):
    add_tb(s, cx+.05, ry+.06, cw-.1, row_h-.1, v, 8.5, col, bold=bld); cx+=cw

add_notes(s,
    f"Commission earnings.\n"
    f"Total pool: {fm(K['comm25'])} → {fm(K['comm26'])} (+{round((K['comm26']-K['comm25'])/K['comm25']*100)}%, +{fm(comm_diff)}).\n"
    f"Effective blended rate: {K['comm_rate26']}% (base 8%, initiative bracket 12%).\n"
    "Top 8 DCs shown. Commission uplift cost is more than offset by revenue and GP growth.")


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — CONCLUSION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 0.55, 7.5, TEAL)

add_tb(s, 0.85, 0.3, 11, .65, "Conclusion & Recommendations", 30, WHITE, bold=True, font="Cambria")
add_tb(s, 0.85, 0.95, 11, .3, "Revenue Management Initiative  ·  H1 2026", 12, RGBColor(0x7A,0xA8,0xCC))

items = [
    ("✓","Initiative is Working",
     f"Bracket adoption grew from {K['init_pct25']}% to {K['init_pct26']}% of $5K+ jobs (+{init_growth}%). "
     f"$5K+ revenue up {rev_growth}% YoY to {fm(K['rev26'])}. Total intake reached {fm(K['all_rev26'])} (+{round((K['all_rev26']-K['all_rev25'])/K['all_rev25']*100,1)}%).",
     GREEN),
    ("✓","Conversion Improving",
     f"$5K+ conversion rate improved {K['conv_over25']}% → {K['conv_over26']}% across {fn(K['acts26'])} sales appointments "
     f"(+{fn(K['acts26']-K['acts25'])} vs 2025). More appointments and a better close rate — lead quality is improving.",
     TEAL),
    ("⚠","Monitor GP% Erosion",
     f"Margin rate declined {abs(K['gp_pct26']-K['gp_pct25']):.1f}pp to {K['gp_pct26']}%. "
     f"GP dollars grew to {fm(K['gp26'])} (+{gp_growth}%). Watch the 40–50% brackets — {D['buckets'][5]['count26']} jobs vs {D['buckets'][5]['count25']} last year.",
     AMBER),
    ("→","Recommendation: Expand & Refine",
     f"Set rep-level adoption targets (suggested: 40%+ of $5K+ jobs). Review bracket at Q3 — "
     f"consider 35–38% at 12% / 38–40% at 10% to recover margin while sustaining momentum. "
     f"Share per-rep data to recognise top performers.",
     BLUE),
]
for i,(icon,title,body,col) in enumerate(items):
    ry = 1.45+i*1.38
    add_rect(s, 0.85, ry, 12.0, 1.25, RGBColor(0x1E,0x3A,0x5F), col, 0.8)
    circ2 = s.shapes.add_shape(9, Inches(0.95), Inches(ry+.2), Inches(.75), Inches(.75))
    circ2.fill.solid(); circ2.fill.fore_color.rgb = col; circ2.line.fill.background()
    add_tb(s, 0.95, ry+.2, .75, .75, icon, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_tb(s, 1.85, ry+.1, 10.8, .32, title, 11, WHITE, bold=True)
    add_tb(s, 1.85, ry+.43, 10.8, .72, body, 9.5, RGBColor(0xCB,0xD5,0xE1))

add_tb(s, 0.85, 7.1, 11.5, .28,
    "Prepared by Sales Analytics  ·  Data as at 27 June 2026  ·  Ex-GST  ·  Confirmed orders  ·  Sales appointments only",
    7.5, SLATE, align=PP_ALIGN.CENTER)

add_notes(s,
    f"Conclusion.\n"
    f"1. Initiative delivering: {K['init_pct26']}% bracket adoption, revenue +{rev_growth}%, total intake {fm(K['all_rev26'])}.\n"
    f"2. Conversion up to {K['conv_over26']}% on $5K+ jobs — both volume and rate improved.\n"
    f"3. GP% watch: -{abs(K['gp_pct26']-K['gp_pct25']):.1f}pp but GP$ grew. Acceptable trade-off currently.\n"
    "4. Actions: rep-level targets, Q3 bracket review, share consultant performance data.")


# ── SAVE ───────────────────────────────────────────────────
out = "Revenue_Initiative_H1_2026_v4.pptx"
prs.save(out)
print(f"Done — {out} ({prs.slides.__len__()} slides)")
