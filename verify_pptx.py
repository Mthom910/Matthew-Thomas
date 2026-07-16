from pptx import Presentation
prs = Presentation("Revenue_Initiative_H1_2026.pptx")
print(f"Slides: {len(prs.slides)}")
print(f"Width: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print()
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t and len(t) > 3:
                texts.append(t[:55].replace("\n"," "))
    notes = slide.notes_slide.notes_text_frame.text[:100].replace("\n"," ")
    print(f"Slide {i}:")
    for t in texts[:4]:
        print(f"  [{t}]")
    print(f"  Notes: {notes[:80]}")
    print()
